"""House style for Deutsche Telekom Business Solutions brochures.

Built with python-pptx (absolute shape positioning, not HTML/CSS flow layout)
following the pattern in agents/4_langchain_langgraph/slide_kit.py. Convert
the resulting .pptx to PDF with pptx_to_pdf.convert_pptx_to_pdf.

This module builds the empty/skeleton template only (placeholder content for
all 9 sections) — wiring real solution/proposal content in is a follow-up.
"""

from datetime import date

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

MAGENTA = RGBColor(0xE2, 0x00, 0x74)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF7, 0xF7, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.7)
CONTENT_WIDTH = SLIDE_WIDTH - 2 * MARGIN

SECTION_TITLES = [
    "Cover",
    "Executive Summary",
    "Customer Challenges",
    "Recommended Solutions",
    "Business Benefits",
    "Implementation Approach",
    "Why Deutsche Telekom",
    "Next Steps",
    "Contact Information",
]


def _new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _textbox(slide, left, top, width, height, text, size, color, bold=False, align=None, spacing=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    para = frame.paragraphs[0]
    if align is not None:
        para.alignment = align
    if spacing is not None:
        para.line_spacing = spacing
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return box


def _rule(slide, left, top, width, height=Pt(2.5), color=MAGENTA):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _bullet_list(slide, left, top, width, items, size=15, color=DARK_GRAY, gap=Inches(0.65), marker_color=MAGENTA):
    for i, item in enumerate(items):
        y = top + i * gap
        marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, y + Inches(0.1), Inches(0.14), Inches(0.14))
        marker.fill.solid()
        marker.fill.fore_color.rgb = marker_color
        marker.line.fill.background()
        _textbox(slide, left + Inches(0.32), y, width - Inches(0.32), gap, item, size, color)


def _section_header(slide, title, kicker):
    _textbox(slide, MARGIN, Inches(0.5), CONTENT_WIDTH, Inches(0.35), kicker.upper(), 11, MAGENTA, bold=True)
    _textbox(slide, MARGIN, Inches(0.85), CONTENT_WIDTH, Inches(0.75), title, 28, DARK_GRAY, bold=True)
    _rule(slide, MARGIN, Inches(1.62), Inches(1.6))


def _footer(slide, page_num, total):
    _textbox(slide, MARGIN, SLIDE_HEIGHT - Inches(0.5), Inches(7), Inches(0.3),
              "Deutsche Telekom Business Solutions", 9, MID_GRAY)
    _textbox(slide, SLIDE_WIDTH - MARGIN - Inches(2), SLIDE_HEIGHT - Inches(0.5), Inches(2), Inches(0.3),
              f"Page {page_num} of {total}", 9, MID_GRAY, align=PP_ALIGN.RIGHT)


def _build_cover_slide(prs, prepared_for: str, date_str: str):
    slide = _blank_slide(prs)
    width = Inches(9)
    left = (SLIDE_WIDTH - width) // 2

    _textbox(slide, left, Inches(2.1), width, Inches(0.4), "DEUTSCHE TELEKOM", 16, MAGENTA,
             bold=True, align=PP_ALIGN.CENTER, spacing=1.0)
    _rule(slide, (SLIDE_WIDTH - Inches(1.8)) // 2, Inches(2.62), Inches(1.8))
    _textbox(slide, left, Inches(2.9), width, Inches(0.9), "Business Solutions", 36, DARK_GRAY,
             align=PP_ALIGN.CENTER)

    _textbox(slide, left, Inches(4.35), width, Inches(0.35), "Prepared for", 14, MID_GRAY, align=PP_ALIGN.CENTER)
    _textbox(slide, left, Inches(4.7), width, Inches(0.5), prepared_for, 20, DARK_GRAY, bold=True,
             align=PP_ALIGN.CENTER)

    _textbox(slide, left, Inches(5.45), width, Inches(0.35), date_str, 12, MID_GRAY, align=PP_ALIGN.CENTER)
    _textbox(slide, left, Inches(6.6), width, Inches(0.35), "CONFIDENTIAL", 10, MID_GRAY,
             align=PP_ALIGN.CENTER, spacing=1.0)


def _build_executive_summary_slide(prs, page_num, total, points):
    slide = _blank_slide(prs)
    _section_header(slide, "Executive Summary", "Overview")
    _bullet_list(slide, MARGIN, Inches(2.1), CONTENT_WIDTH, points, size=16, gap=Inches(0.75))
    _footer(slide, page_num, total)


def _build_customer_challenges_slide(prs, page_num, total, challenges):
    slide = _blank_slide(prs)
    _section_header(slide, "Customer Challenges", "Where things stand today")
    _bullet_list(slide, MARGIN, Inches(2.1), CONTENT_WIDTH, challenges, size=16, gap=Inches(0.75))
    _footer(slide, page_num, total)


def _solution_block(slide, left, top, width, height, name, category, description):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_GRAY
    bg.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.07), height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = MAGENTA
    accent.line.fill.background()
    _textbox(slide, left + Inches(0.3), top + Inches(0.18), width - Inches(0.6), Inches(0.4), name, 17, MAGENTA, bold=True)
    _textbox(slide, left + Inches(0.3), top + Inches(0.62), width - Inches(0.6), Inches(0.3), category.upper(), 10, MID_GRAY)
    _textbox(slide, left + Inches(0.3), top + Inches(1.0), width - Inches(0.6), height - Inches(1.2), description, 13, DARK_GRAY)


def _build_recommended_solutions_slide(prs, page_num, total, solutions):
    slide = _blank_slide(prs)
    _section_header(slide, "Recommended Solutions", "What we propose")
    block_height = Inches(1.8)
    gap = Inches(0.25)
    y = Inches(2.1)
    for name, category, description in solutions:
        _solution_block(slide, MARGIN, y, CONTENT_WIDTH, block_height, name, category, description)
        y += block_height + gap
    _footer(slide, page_num, total)


def _build_business_benefits_slide(prs, page_num, total, benefits):
    slide = _blank_slide(prs)
    _section_header(slide, "Business Benefits", "What this means for you")
    _bullet_list(slide, MARGIN, Inches(2.1), CONTENT_WIDTH, benefits, size=16, gap=Inches(0.75))
    _footer(slide, page_num, total)


def _build_implementation_approach_slide(prs, page_num, total, phases):
    slide = _blank_slide(prs)
    _section_header(slide, "Implementation Approach", "How we get there")
    phase_width = (CONTENT_WIDTH - Inches(0.6)) // 3
    x = MARGIN
    for i, (phase_name, phase_desc) in enumerate(phases[:3], start=1):
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(2.2), Inches(0.5), Inches(0.5))
        badge.fill.solid()
        badge.fill.fore_color.rgb = MAGENTA
        badge.line.fill.background()
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = badge.text_frame.paragraphs[0].add_run()
        run.text = str(i)
        run.font.color.rgb = WHITE
        run.font.bold = True
        run.font.size = Pt(16)

        _textbox(slide, x, Inches(2.9), phase_width, Inches(0.4), phase_name, 15, DARK_GRAY, bold=True)
        _textbox(slide, x, Inches(3.35), phase_width, Inches(2), phase_desc, 12, MID_GRAY)
        x += phase_width + Inches(0.3)
    _footer(slide, page_num, total)


def _build_why_dt_slide(prs, page_num, total, reasons):
    slide = _blank_slide(prs)
    _section_header(slide, "Why Deutsche Telekom", "A partner you can rely on")
    _bullet_list(slide, MARGIN, Inches(2.1), CONTENT_WIDTH, reasons, size=16, gap=Inches(0.75))
    _footer(slide, page_num, total)


def _build_next_steps_slide(prs, page_num, total, steps):
    slide = _blank_slide(prs)
    _section_header(slide, "Next Steps", "Let's move forward")
    _bullet_list(slide, MARGIN, Inches(2.1), CONTENT_WIDTH, steps, size=16, gap=Inches(0.75))
    _footer(slide, page_num, total)


def _build_contact_slide(prs, contact_line: str):
    slide = _blank_slide(prs)
    width = Inches(9)
    left = (SLIDE_WIDTH - width) // 2

    _textbox(slide, left, Inches(2.6), width, Inches(0.6), "Let's talk.", 30, DARK_GRAY, align=PP_ALIGN.CENTER)
    _rule(slide, (SLIDE_WIDTH - Inches(1.8)) // 2, Inches(3.35), Inches(1.8))
    _textbox(slide, left, Inches(3.7), width, Inches(0.4), "Need help?", 14, MID_GRAY, align=PP_ALIGN.CENTER)
    _textbox(slide, left, Inches(4.1), width, Inches(0.4), contact_line, 16, DARK_GRAY,
             bold=True, align=PP_ALIGN.CENTER)

    _textbox(slide, left, Inches(6.6), width, Inches(0.4), "DEUTSCHE TELEKOM", 13, MAGENTA,
             bold=True, align=PP_ALIGN.CENTER, spacing=1.0)


def build_brochure_pptx(outfile: str) -> None:
    """Build the empty/skeleton brochure with placeholder content for all 9 sections."""
    prs = _new_presentation()
    total = len(SECTION_TITLES)

    _build_cover_slide(prs, prepared_for="[Customer Company Name]", date_str=date.today().strftime("%d %B %Y"))

    _build_executive_summary_slide(prs, 2, total, [
        "[One-sentence summary of the customer's situation and opportunity]",
        "[Why this matters to them right now]",
        "[What Deutsche Telekom is proposing, in one line]",
    ])

    _build_customer_challenges_slide(prs, 3, total, [
        "Challenge 1: [pain point currently affecting the business]",
        "Challenge 2: [pain point currently affecting the business]",
        "Challenge 3: [pain point currently affecting the business]",
    ])

    _build_recommended_solutions_slide(prs, 4, total, [
        ("[Solution Name]", "[Category]", "[One or two sentence description of what it does and why it fits.]"),
        ("[Solution Name]", "[Category]", "[One or two sentence description of what it does and why it fits.]"),
    ])

    _build_business_benefits_slide(prs, 5, total, [
        "[Quantifiable or concrete benefit #1]",
        "[Quantifiable or concrete benefit #2]",
        "[Quantifiable or concrete benefit #3]",
    ])

    _build_implementation_approach_slide(prs, 6, total, [
        ("Discovery", "[What happens in this phase]"),
        ("Pilot", "[What happens in this phase]"),
        ("Rollout", "[What happens in this phase]"),
    ])

    _build_why_dt_slide(prs, 7, total, [
        "[Credibility point — scale, experience, certifications]",
        "[Credibility point — support, reliability]",
        "[Credibility point — compliance, data residency]",
    ])

    _build_next_steps_slide(prs, 8, total, [
        "[Immediate next action, e.g. a scoping call]",
        "[What Deutsche Telekom will prepare]",
        "[What we need from the customer]",
    ])

    _build_contact_slide(prs, contact_line="[Contact details to be confirmed]")

    prs.save(outfile)
